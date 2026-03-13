#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT 页面快速定位工具（增强版）
功能：
1. 构建索引
2. 搜索内容
3. 提取完整页面
4. 批量提取匹配页面
5. 列出所有内容
6. 打开指定页面
"""
import os
import json
import sys
from pptx import Presentation
from pathlib import Path

def build_index(ppt_dir, index_file):
    """构建 PPT 索引"""
    index = []
    
    print("=" * 70)
    print("PPT 页面索引构建工具")
    print("=" * 70)
    print(f"\nScanning: {ppt_dir}")
    
    for filename in sorted(os.listdir(ppt_dir)):
        if not filename.lower().endswith(('.ppt', '.pptx')):
            continue
        
        ppt_path = os.path.join(ppt_dir, filename)
        
        try:
            prs = Presentation(ppt_path)
            slides = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                text_parts = []
                shapes_count = 0
                images_count = 0
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                    
                    # 统计形状
                    shapes_count += 1
                    
                    # 统计图片
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        images_count += 1
                
                if text_parts:
                    page_text = "\n".join(text_parts)
                    
                    # 提取标题（通常是第一个文本框的前几个词）
                    title = text_parts[0][:50] if text_parts else f"Page {slide_idx}"
                    
                    slides.append({
                        "page": slide_idx,
                        "title": title,
                        "text_preview": page_text[:300] + "..." if len(page_text) > 300 else page_text,
                        "shapes_count": shapes_count,
                        "images_count": images_count
                    })
            
            if slides:
                index.append({
                    "file": filename,
                    "path": ppt_path,
                    "total_slides": len(slides),
                    "slides": slides
                })
            
            print(f"✓ {filename} ({len(slides)} pages)")
            
        except Exception as e:
            print(f"✗ {filename} - {e}")
    
    # Save index
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(index, f, ensure_ascii=False, indent=2)
    
    total_pages = sum(len(item['slides']) for item in index)
    
    print("\n" + "=" * 70)
    print("Index built successfully!")
    print("=" * 70)
    print(f"Total files: {len(index)}")
    print(f"Total pages: {total_pages}")
    print(f"Index file: {index_file}")

def search(index_file, query, max_results=20):
    """搜索 PPT 内容"""
    
    if not os.path.exists(index_file):
        print(f"Index file not found: {index_file}")
        print("Please run: python ppt_tool.py --build")
        return []
    
    # Load index
    with open(index_file, 'r', encoding='utf-8') as f:
        index = json.load(f)
    
    print("=" * 70)
    print("PPT Content Search")
    print("=" * 70)
    print(f"Query: {query}")
    print(f"Indexed files: {len(index)}")
    
    query_lower = query.lower()
    results = []
    
    # Search each page
    for file_entry in index:
        filename = file_entry['file']
        path = file_entry['path']
        
        for slide in file_entry['slides']:
            page_num = slide['page']
            text = slide['text_preview']
            title = slide['title']
            
            # Check if matches
            if query_lower in text.lower() or query_lower in title.lower():
                results.append({
                    "filename": filename,
                    "path": path,
                    "page": page_num,
                    "title": title,
                    "text": text
                })
    
    # Show results
    print(f"\nFound {len(results)} matching results")
    print("=" * 70)
    
    if results:
        for i, result in enumerate(results[:max_results], 1):
            print(f"\n[{i}] File: {result['filename']}")
            print(f"    Page: {result['page']}")
            print(f"    Title: {result['title']}")
            print(f"    Path: {result['path']}")
            
            # Show preview
            text = result['text']
            idx = text.lower().find(query_lower) if query_lower in text.lower() else -1
            
            if idx >= 0:
                start = max(0, idx - 80)
                end = min(len(text), idx + len(query) + 80)
                
                preview = text[start:end]
                if start > 0:
                    preview = "..." + preview
                if end < len(text):
                    preview = preview + "..."
                
                # Highlight match
                preview = preview.replace(query, f"【{query}】")
                print(f"    Preview: {preview}")
            else:
                print(f"    Preview: {text[:150]}...")
        
        if len(results) > max_results:
            print(f"\n... and {len(results) - max_results} more results")
    else:
        print("\nNo matches found")
    
    return results

def extract_page(ppt_path, page_num, output_file=None):
    """提取 PPT 指定页面的完整内容"""
    try:
        prs = Presentation(ppt_path)
        
        # Check if page exists
        if page_num < 1 or page_num > len(prs.slides):
            print(f"✗ Invalid page number: {page_num}")
            print(f"   Valid range: 1-{len(prs.slides)}")
            return None
        
        slide = prs.slides[page_num - 1]
        
        content = {
            "filename": os.path.basename(ppt_path),
            "path": ppt_path,
            "page": page_num,
            "title": "",
            "text_blocks": [],
            "full_text": "",
            "images_count": 0,
            "shapes_count": 0
        }
        
        for shape in slide.shapes:
            content['shapes_count'] += 1
            
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                content['text_blocks'].append(text)
                
                if not content['title']:
                    content['title'] = text[:50]
            
            # Count images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                content['images_count'] += 1
        
        content['full_text'] = "\n\n".join(content['text_blocks'])
        
        # Display content
        print("=" * 70)
        print(f"Page {page_num} Extracted: {content['filename']}")
        print("=" * 70)
        print(f"File: {ppt_path}")
        print(f"Shapes: {content['shapes_count']}")
        print(f"Images: {content['images_count']}")
        print(f"Title: {content['title']}")
        print("\n--- Full Text ---\n")
        print(content['full_text'])
        
        # Save to file if requested
        if output_file:
            save_content(content, output_file)
            print(f"\n✓ Saved to: {output_file}")
        
        return content
        
    except Exception as e:
        print(f"✗ Error extracting page: {e}")
        return None

def extract_all_matching(index_file, query, output_dir=None):
    """提取所有匹配页面的完整内容"""
    
    if not os.path.exists(index_file):
        print(f"Index file not found: {index_file}")
        print("Please run: python ppt_tool.py --build")
        return []
    
    # Search first
    results = search(index_file, query, max_results=1000)
    
    if not results:
        print("\nNo results to extract")
        return []
    
    print("\n" + "=" * 70)
    print("Extracting All Matching Pages")
    print("=" * 70)
    
    extracted_contents = []
    
    # Create output directory if specified
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
        print(f"Output directory: {output_dir}\n")
    
    for i, result in enumerate(results, 1):
        print(f"\n[{i}/{len(results)}] Extracting: {result['filename']} - Page {result['page']}")
        
        content = extract_page(
            result['path'],
            result['page'],
            output_file=os.path.join(output_dir, f"{i:02d}_{result['filename']}_page{result['page']}.txt") if output_dir else None
        )
        
        if content:
            extracted_contents.append(content)
        
        print()  # Empty line between results
    
    print("=" * 70)
    print(f"✓ Extracted {len(extracted_contents)} pages")
    print("=" * 70)
    
    return extracted_contents

def save_content(content, output_file):
    """保存提取的内容到文件"""
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write("=" * 70 + "\n")
        f.write(f"PPT Page Extract: {content['filename']}\n")
        f.write("=" * 70 + "\n\n")
        f.write(f"Page: {content['page']}\n")
        f.write(f"Path: {content['path']}\n")
        f.write(f"Shapes: {content['shapes_count']}\n")
        f.write(f"Images: {content['images_count']}\n")
        f.write(f"Title: {content['title']}\n")
        f.write("\n" + "-" * 70 + "\n")
        f.write("Full Text:\n")
        f.write("-" * 70 + "\n\n")
        f.write(content['full_text'] + "\n")

def list_content(index_file, max_pages=50):
    """列出所有 PPT 内容"""
    
    if not os.path.exists(index_file):
        print(f"Index file not found: {index_file}")
        print("Please run: python ppt_tool.py --build")
        return
    
    # Load index
    with open(index_file, 'r', encoding='utf-8') as f:
        index = json.load(f)
    
    print("=" * 70)
    print("All PPT Content")
    print("=" * 70)
    
    for file_entry in index:
        filename = file_entry['file']
        print(f"\n📄 {filename}")
        print("-" * 70)
        
        for slide in file_entry['slides'][:max_pages]:
            page_num = slide['page']
            title = slide['title']
            text = slide['text_preview']
            
            print(f"\nPage {page_num}: {title}")
            print(f"Text: {text[:200]}..." if len(text) > 200 else text)

def open_ppt(ppt_path, page_num=None):
    """打开 PPT 文件"""
    import subprocess
    
    try:
        subprocess.Popen(['start', '', ppt_path], shell=True)
        print(f"✓ Opened: {ppt_path}")
        
        if page_num:
            print(f"  Please navigate to page {page_num}")
    except Exception as e:
        print(f"✗ Error opening file: {e}")
        print(f"  Please open manually: {ppt_path}")

def load_config(config_file):
    """加载配置文件"""
    try:
        with open(config_file, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"Error loading config file: {e}")
        # 返回默认配置
        return {
            "ppt_dir": r"C:\Users\Administrator\Desktop\PPT集合",
            "index_file": r"C:\Users\Administrator\AppData\Roaming\memu-bot\agent-output\ppt_index.json",
            "default_output_dir": r"C:\Users\Administrator\AppData\Roaming\memu-bot\agent-output\ppt_extracted"
        }

def main():
    # 加载配置
    config_file = os.path.join(os.path.dirname(__file__), 'ppt_config.json')
    config = load_config(config_file)
    ppt_dir = config.get("ppt_dir")
    index_file = config.get("index_file")
    default_output_dir = config.get("default_output_dir")
    
    if len(sys.argv) < 2:
        print("=" * 70)
        print("PPT Page Finder - Enhanced Tool")
        print("=" * 70)
        print("\nUsage:")
        print("  python ppt_tool.py --build                           # Build index")
        print("  python ppt_tool.py --search <query>                   # Search content")
        print("  python ppt_tool.py --extract <filename> <page>        # Extract single page")
        print("  python ppt_tool.py --extract-all <query>               # Extract all matching pages")
        print("  python ppt_tool.py --extract-all <query> <output-dir> # Extract to custom directory")
        print("  python ppt_tool.py --list                            # List all content")
        print("  python ppt_tool.py --open <filename> [page]           # Open PPT file")
        print("\nExamples:")
        print("  python ppt_tool.py --build")
        print("  python ppt_tool.py --search AI")
        print("  python ppt_tool.py --search 金融分析")
        print("  python ppt_tool.py --extract ai-agent-tech-insights.pptx 1")
        print("  python ppt_tool.py --extract-all AI")
        print("  python ppt_tool.py --extract-all OpenClaw ./extracted")
        print("  python ppt_tool.py --open ai-agent-tech-insights.pptx")
        return
    
    command = sys.argv[1]
    
    if command == "--build":
        build_index(ppt_dir, index_file)
    
    elif command == "--search":
        if len(sys.argv) < 3:
            print("Please provide search query")
            print("Usage: python ppt_tool.py --search <query>")
            return
        query = " ".join(sys.argv[2:])
        search(index_file, query)
    
    elif command == "--extract":
        if len(sys.argv) < 4:
            print("Please provide filename and page number")
            print("Usage: python ppt_tool.py --extract <filename> <page>")
            return
        
        filename = sys.argv[2]
        page_num = int(sys.argv[3])
        
        # Find file in index
        if os.path.exists(index_file):
            with open(index_file, 'r', encoding='utf-8') as f:
                index = json.load(f)
            
            for file_entry in index:
                if file_entry['file'].lower() == filename.lower():
                    extract_page(file_entry['path'], page_num)
                    return
            
            print(f"✗ File not found in index: {filename}")
        else:
            print(f"Index file not found: {index_file}")
            print("Please run: python ppt_tool.py --build")
    
    elif command == "--extract-all":
        if len(sys.argv) < 3:
            print("Please provide search query")
            print("Usage: python ppt_tool.py --extract-all <query> [output-dir]")
            return
        
        query = sys.argv[2]
        output_dir = sys.argv[3] if len(sys.argv) > 3 else default_output_dir
        
        extract_all_matching(index_file, query, output_dir)
    
    elif command == "--list":
        list_content(index_file)
    
    elif command == "--open":
        if len(sys.argv) < 3:
            print("Please provide filename")
            print("Usage: python ppt_tool.py --open <filename> [page]")
            return
        
        filename = sys.argv[2]
        page_num = int(sys.argv[3]) if len(sys.argv) > 3 else None
        
        # Find file in index
        if os.path.exists(index_file):
            with open(index_file, 'r', encoding='utf-8') as f:
                index = json.load(f)
            
            for file_entry in index:
                if file_entry['file'].lower() == filename.lower():
                    open_ppt(file_entry['path'], page_num)
                    return
            
            print(f"✗ File not found in index: {filename}")
        else:
            print(f"Index file not found: {index_file}")
            print("Please run: python ppt_tool.py --build")
    
    else:
        print(f"Unknown command: {command}")
        print("Available commands: --build, --search, --extract, --extract-all, --list, --open")

if __name__ == '__main__':
    main()
